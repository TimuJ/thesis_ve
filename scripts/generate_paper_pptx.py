"""
Generate a PowerPoint presentation for arxiv 2505.20171:
"Long-Context State-Space Video World Models" (Po et al., 2025)

Focus: method, technical details, innovations, what made it publishable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)  # deep navy
BG_MED    = RGBColor(0x22, 0x22, 0x3A)
ACCENT    = RGBColor(0x4E, 0xC9, 0xB0)  # teal accent
ACCENT2   = RGBColor(0x6C, 0x9B, 0xD2)  # soft blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE    = RGBColor(0xE8, 0x8D, 0x3F)
RED_SOFT  = RGBColor(0xE8, 0x6B, 0x6B)
GREEN     = RGBColor(0x6B, 0xE8, 0x6B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items, font_size=18,
                     color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_number(slide, number, top=Inches(0.6)):
    add_textbox(slide, Inches(0.8), top, Inches(1), Inches(0.6),
                f"{number:02d}", font_size=36, color=ACCENT, bold=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(4), Pt(4))

add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.5),
            "Long-Context State-Space\nVideo World Models",
            font_size=40, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
            "Ryan Po, Yotam Nitzan, Richard Zhang, Berlin Chen, Tri Dao,\n"
            "Eli Shechtman, Gordon Wetzstein, Xun Huang",
            font_size=18, color=LIGHT)
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
            "Stanford University  |  Princeton University  |  Adobe Research",
            font_size=16, color=ACCENT)
add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(0.4),
            "arXiv 2505.20171  |  CVPR 2025 area  |  cs.CV",
            font_size=14, color=LIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Problem Statement
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 1)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "The Problem: Video World Models Forget", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "Video diffusion models generate frames autoregressively via sliding-window attention",
    "Sliding window = limited context = NO long-term memory",
    "Result: environment completely changes when agent looks away and back",
    "Example: in a game, looking right then left → entire scene regenerated differently",
    "",
    "Why can't we just extend the attention window?",
    "  (1) Training cost scales QUADRATICALLY with context length",
    "  (2) Per-frame inference time grows LINEARLY — too slow for real-time",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.0),
                 items, font_size=20, color=LIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Key Innovation (high-level)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 2)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Key Innovation: SSMs for Causal World State Tracking",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "Replace attention's temporal modeling with State-Space Models (SSMs / Mamba)",
    "",
    "Why SSMs?",
    "  • Causal by nature — perfect for autoregressive frame generation",
    "  • Fixed-size hidden state compresses entire history → O(1) memory at inference",
    "  • Linear training complexity vs. quadratic for transformers",
    "",
    "Critical insight: previous SSM-for-vision works use BIDIRECTIONAL scans",
    "  → breaks causality, can't do efficient AR inference",
    "This paper: UNIDIRECTIONAL SSM for temporal dynamics — exploits SSMs' true strength",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.0),
                 items, font_size=20, color=LIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Complexity Comparison Table
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 3)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Complexity Comparison: Why This Matters",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Table header
headers = ["Architecture", "Training", "AR Inference\n(total)", "AR Inference\n(per frame)", "Long\nMemory"]
col_lefts = [Inches(1.5), Inches(4.5), Inches(6.5), Inches(8.5), Inches(10.5)]
col_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(1.5)]

for i, h in enumerate(headers):
    add_textbox(slide, col_lefts[i], Inches(1.8), col_widths[i], Inches(0.7),
                h, font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ("Bidirectional Attention", "Quadratic", "Cubic", "Quadratic", "Yes"),
    ("Causal Attention", "Quadratic", "Quadratic", "Linear", "Yes"),
    ("Causal + Sliding Window", "Sub-quadratic", "Linear", "Constant", "No"),
    ("Ours (SSM + Local Attn)", "Linear", "Linear", "Constant", "Yes"),
]

for row_idx, row in enumerate(rows):
    y = Inches(2.6 + row_idx * 0.7)
    is_ours = row_idx == 3
    for col_idx, val in enumerate(row):
        c = ACCENT if is_ours else LIGHT
        b = is_ours
        if col_idx == 4:
            if val == "Yes" and is_ours:
                c = GREEN
            elif val == "No":
                c = RED_SOFT
        add_textbox(slide, col_lefts[col_idx], y, col_widths[col_idx], Inches(0.5),
                    val, font_size=16, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# highlight bar behind "Ours" row
rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1.3), Inches(2.6 + 3 * 0.7 - Inches(0.05).emu / 914400 * 914400),
    Inches(11.0), Inches(0.55))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0x2A, 0x3A, 0x2A)
rect.line.fill.background()
# send to back by reordering
sp = rect._element
sp.getparent().remove(sp)
slide.shapes._spTree.insert(2, sp)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Architecture Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 4)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Architecture: Two Core Components",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Left box — Block-wise SSM
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.0), Inches(1.8), Inches(5.5), Inches(4.8))
box1.fill.solid()
box1.fill.fore_color.rgb = BG_MED
box1.line.color.rgb = ACCENT
box1.line.width = Pt(2)

add_textbox(slide, Inches(1.3), Inches(2.0), Inches(5.0), Inches(0.5),
            "1. Block-wise SSM Scan", font_size=24, color=ACCENT, bold=True)

items_left = [
    "Divide spatial dims into blocks of (bh, bw, T)",
    "Each block gets its OWN independent Mamba scan",
    "Temporal tokens now separated by bh×bw (not H×W)",
    "→ Much closer temporal neighbors for SSM",
    "",
    "Trade-off controlled by block size:",
    "  Small blocks → better temporal memory",
    "  Large blocks → better spatial coherence",
    "  Solution: vary (bh, bw) across layers",
    "",
    "Bonus: effectively increases SSM state dim",
    "(each block has separate state → more capacity)",
]
add_bullet_slide(slide, Inches(1.3), Inches(2.6), Inches(5.0), Inches(3.8),
                 items_left, font_size=15, color=LIGHT)

# Right box — Frame Local Attention
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8))
box2.fill.solid()
box2.fill.fore_color.rgb = BG_MED
box2.line.color.rgb = ACCENT2
box2.line.width = Pt(2)

add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.0), Inches(0.5),
            "2. Frame Local Attention", font_size=24, color=ACCENT2, bold=True)

items_right = [
    "After every Mamba layer → local attention block",
    "Block-wise causal attention over k previous frames",
    "",
    "Attention mask: M[i,j] = 1 if j ∈ [i-k, i]",
    "Bidirectional WITHIN frames",
    "Causal ACROSS frames (window of k)",
    "",
    "Why needed?",
    "  Mamba struggles with associative recall",
    "  Local attention fixes frame-wise quality",
    "  Fixes short-term temporal consistency",
    "",
    "Result: hybrid SSM+Attention architecture",
]
add_bullet_slide(slide, Inches(7.1), Inches(2.6), Inches(5.0), Inches(3.8),
                 items_right, font_size=15, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Block-wise Scan Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 5)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Block-wise SSM Scan: The Core Technical Detail",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "Standard approach: flatten all spatiotemporal tokens → single scan",
    "  Spatial-major ordering: (h1,w1,t1), (h1,w2,t1), ... (hH,wW,t1), (h1,w1,t2), ...",
    "  Problem: temporal neighbors are H×W tokens apart in scan order",
    "  → SSM \"forgets\" temporal info across such large gaps",
    "",
    "This paper's approach: divide (H, W) into blocks of (bh, bw)",
    "  Each block scanned independently: (b_h1, b_w1, t1), (b_h1, b_w2, t1), ..., (b_h, b_w, tT)",
    "  Now temporal gap = bh × bw  (e.g., 2×2 = 4 tokens, not 16×16 = 256)",
    "",
    "Layer-dependent block sizes:",
    "  Early layers: smaller blocks → prioritize temporal memory",
    "  Later layers: larger blocks → prioritize spatial coherence",
    "  This multi-scale design is key to the method working",
    "",
    "SSM state dimensionality scales with number of blocks:",
    "  Total state = (H/bh × W/bw) × state_dim_per_block",
    "  → Smaller blocks = more blocks = more total state capacity",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.5),
                 items, font_size=18, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Long-Context Training
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 6)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Training Innovation: Forcing Long-Range Dependencies",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Left: Standard diffusion forcing
add_textbox(slide, Inches(1.5), Inches(1.7), Inches(5), Inches(0.5),
            "Standard Diffusion Forcing", font_size=22, color=RED_SOFT, bold=True)
items_std = [
    "Independent noise levels per frame: ti ~ U(0, T)",
    "ALL frames get noise during training",
    "Model learns to denoise each frame",
    "",
    "Problem: nearby noisy frames are still more useful",
    "than distant noisy frames → model ignores distant context",
    "→ Trapped in LOCAL MINIMUM",
    "→ Never learns long-term dependencies",
]
add_bullet_slide(slide, Inches(1.5), Inches(2.3), Inches(5.0), Inches(4.5),
                 items_std, font_size=16, color=LIGHT)

# Right: Their method
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
            "Their Training Scheme", font_size=22, color=GREEN, bold=True)
items_new = [
    "Keep a random-length PREFIX completely clean (ti = 0)",
    "Add independent noise only to LATER frames",
    "Compute loss only on noised frames",
    "",
    "Why this works:",
    "When later frames have HIGH noise, the clean prefix",
    "is MORE informative than noisy local neighbors",
    "→ Forces model to attend to distant clean context",
    "→ Learns long-range temporal correlations",
    "",
    "Mixed with standard diffusion forcing during training",
]
add_bullet_slide(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
                 items_new, font_size=16, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Efficient Inference
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 7)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Efficient Inference: Constant Cost per Frame",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "At inference, each layer maintains only:",
    "  (1) Fixed-length KV-cache for previous k frames (local attention)",
    "  (2) SSM hidden state per block (constant size, compresses ALL history)",
    "",
    "Memory usage: CONSTANT regardless of video length",
    "  Causal transformer: KV-cache grows linearly with every generated frame",
    "  This method: fixed KV (k frames) + fixed SSM states",
    "",
    "Per-frame generation speed: CONSTANT",
    "  Local attention: only over k frames (not all history)",
    "  SSM update: single state transition per block",
    "",
    "→ Can generate INFINITELY long videos without degradation",
    "→ Suitable for interactive applications (gaming, robotics)",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.0),
                 items, font_size=19, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — Experimental Setup
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 8)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Evaluation: Novel Memory-Focused Tasks",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Datasets
add_textbox(slide, Inches(1.5), Inches(1.7), Inches(5), Inches(0.4),
            "Datasets", font_size=22, color=ACCENT, bold=True)
items_data = [
    "Memory Maze — 3D mazes, 2000 frames/trajectory",
    "  Specifically designed for long-term memory eval",
    "TECO Minecraft — 200K gameplay trajectories",
    "  4 discrete actions, 150 frames each",
]
add_bullet_slide(slide, Inches(1.5), Inches(2.2), Inches(5.0), Inches(2.5),
                 items_data, font_size=17, color=LIGHT)

# Tasks
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
            "Novel Evaluation Tasks", font_size=22, color=ACCENT, bold=True)
items_tasks = [
    "Spatial Retrieval — backtrack through maze,",
    "  generated frames must match original observations",
    "  Tests: can the model RECALL what it saw?",
    "",
    "Spatial Reasoning — continue trajectory forward,",
    "  must reconstruct previously observed regions",
    "  Tests: can the model INFER from memory?",
]
add_bullet_slide(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.0),
                 items_tasks, font_size=17, color=LIGHT)

# Metrics
add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.4),
            "Metrics: SSIM, LPIPS, PSNR — comparing generated vs ground-truth frames",
            font_size=18, color=ORANGE, bold=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — Results
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 9)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Results: Outperforms All Sub-Quadratic Methods",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Retrieval results
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5.5), Inches(0.4),
            "Memory Maze — Retrieval (400 frames)", font_size=20, color=ACCENT, bold=True)

r_headers = ["Model", "SSIM ↑", "LPIPS ↓", "PSNR ↑"]
r_col_lefts = [Inches(1.0), Inches(3.5), Inches(4.8), Inches(6.0)]
r_col_widths = [Inches(2.5), Inches(1.3), Inches(1.2), Inches(1.2)]

for i, h in enumerate(r_headers):
    add_textbox(slide, r_col_lefts[i], Inches(2.2), r_col_widths[i], Inches(0.4),
                h, font_size=14, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

r_rows = [
    ("Causal (192 ctx)", "0.829", "0.147", "26.4", False),
    ("Mamba2", "0.747", "0.313", "20.4", False),
    ("Mamba2 + Local Attn", "0.735", "0.336", "19.3", False),
    ("Ours", "0.898", "0.069", "30.8", True),
    ("Causal (Full ctx)", "0.914", "0.057", "32.6", False),
]

for row_idx, (name, ssim, lpips, psnr, highlight) in enumerate(r_rows):
    y = Inches(2.65 + row_idx * 0.42)
    vals = [name, ssim, lpips, psnr]
    for col_idx, val in enumerate(vals):
        c = GREEN if highlight else (LIGHT if row_idx < 4 else RGBColor(0x88, 0x88, 0x88))
        b = highlight
        add_textbox(slide, r_col_lefts[col_idx], y, r_col_widths[col_idx], Inches(0.35),
                    val, font_size=14, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# Reasoning results
add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.4),
            "Memory Maze — Reasoning (224 frames)", font_size=20, color=ACCENT2, bold=True)

r2_col_lefts = [Inches(7.2), Inches(9.7), Inches(10.9), Inches(12.0)]

for i, h in enumerate(r_headers):
    add_textbox(slide, r2_col_lefts[i], Inches(2.2), r_col_widths[i], Inches(0.4),
                h, font_size=14, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

r2_rows = [
    ("Causal (192 ctx)", "0.839", "0.125", "27.1", False),
    ("Mamba2", "0.827", "0.150", "26.4", False),
    ("Mamba2 + Local Attn", "0.845", "0.113", "27.5", False),
    ("Ours", "0.855", "0.099", "28.2", True),
    ("Causal (Full ctx)", "0.860", "0.089", "28.8", False),
]

for row_idx, (name, ssim, lpips, psnr, highlight) in enumerate(r2_rows):
    y = Inches(2.65 + row_idx * 0.42)
    vals = [name, ssim, lpips, psnr]
    for col_idx, val in enumerate(vals):
        c = GREEN if highlight else (LIGHT if row_idx < 4 else RGBColor(0x88, 0x88, 0x88))
        b = highlight
        add_textbox(slide, r2_col_lefts[col_idx], y, r_col_widths[col_idx], Inches(0.35),
                    val, font_size=14, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# Key takeaway
add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11.5), Inches(0.8),
            "Key: Ours approaches full-context causal transformer quality\n"
            "with LINEAR training cost and CONSTANT inference cost",
            font_size=20, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — Ablations
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 10)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Ablations: Every Component Is Essential",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

# Ablation table
a_headers = ["Configuration", "SSIM ↑", "LPIPS ↓", "PSNR ↑"]
a_col_lefts = [Inches(2.0), Inches(6.5), Inches(8.2), Inches(9.8)]
a_col_widths = [Inches(4.5), Inches(1.5), Inches(1.5), Inches(1.5)]

for i, h in enumerate(a_headers):
    add_textbox(slide, a_col_lefts[i], Inches(1.8), a_col_widths[i], Inches(0.4),
                h, font_size=17, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

a_rows = [
    ("w/o block-wise scan (single scan)", "0.845", "0.113", "27.5"),
    ("Block size = 1 (no spatial coherence)", "0.766", "0.198", "23.1"),
    ("w/o clean-prefix training (Sec 4.2)", "0.809", "0.143", "25.3"),
    ("Full method", "0.855", "0.099", "28.2"),
]

for row_idx, (name, ssim, lpips, psnr) in enumerate(a_rows):
    y = Inches(2.4 + row_idx * 0.55)
    is_full = row_idx == 3
    vals = [name, ssim, lpips, psnr]
    for col_idx, val in enumerate(vals):
        c = GREEN if is_full else LIGHT
        b = is_full
        add_textbox(slide, a_col_lefts[col_idx], y, a_col_widths[col_idx], Inches(0.4),
                    val, font_size=16, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# Interpretation
add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.4),
            "Takeaways:", font_size=22, color=ACCENT, bold=True)
items_abl = [
    "No block-wise scan → temporal tokens too far apart, SSM forgets",
    "Block size 1 → perfect temporal proximity but no spatial coherence → bad reasoning",
    "No clean-prefix training → model falls into local minimum, ignores distant context",
    "All three innovations are NECESSARY — removing any one significantly hurts performance",
]
add_bullet_slide(slide, Inches(1.5), Inches(5.3), Inches(10.5), Inches(2.0),
                 items_abl, font_size=17, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — What Made This Publishable
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 11)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Why This Paper Gets Published",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "1. CLEAR PROBLEM FRAMING",
    "   Existing video world models forget → concretely demonstrated (Fig. 1: look right, look left)",
    "",
    "2. PRINCIPLED SOLUTION — not just engineering",
    "   SSMs are naturally causal → used for their INTENDED purpose (unlike prior bidirectional hacks)",
    "   Block-wise scan: theoretically motivated trade-off with clear complexity analysis",
    "",
    "3. THREE SYNERGISTIC CONTRIBUTIONS, each ablated",
    "   Architecture (block-wise SSM + local attention)",
    "   Training (clean prefix forcing)",
    "   Inference (constant-cost via fixed state)",
    "",
    "4. NOVEL EVALUATION PROTOCOL",
    "   Spatial retrieval + spatial reasoning tasks — specifically test long-term memory",
    "   Not just FID/FVD but task-specific metrics that validate the core claim",
    "",
    "5. STRONG RESULTS + HONEST POSITIONING",
    "   Approaches full-context transformer (the expensive upper bound) at linear cost",
    "   Doesn't claim to beat everything — clearly positions as sub-quadratic champion",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.5),
                 items, font_size=16, color=LIGHT)

# Color the numbered items
# (already readable as-is with bullet formatting)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — Limitations & Future Work
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, 12)
add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.6),
            "Limitations & Future Directions",
            font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1.5), Inches(1.25), Inches(3), Pt(3))

items = [
    "Limitations (acknowledged by authors):",
    "  • Not yet interactive frame rates — needs timestep distillation",
    "  • Cannot extrapolate beyond training context length",
    "  • Experiments limited to low-resolution synthetic videos",
    "",
    "Future directions:",
    "  • Timestep distillation for real-time generation",
    "  • Length extrapolation techniques (DeciMamba, etc.)",
    "  • Scaling to high-resolution realistic videos",
    "",
    "Connection to VSR (our thesis direction):",
    "  • SSM temporal backbone → natural fit for long-video SR (>1 min)",
    "  • Block-wise scan → could propagate SR features across long sequences",
    "  • Diffusion + SSM hybrid → high-quality SR with efficient temporal modeling",
    "  • Clean-prefix training → could help SR models leverage distant reference frames",
]
add_bullet_slide(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(5.5),
                 items, font_size=18, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — Summary
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "Summary", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(3), Pt(3))

items = [
    "Problem: Video world models lack long-term memory due to attention's O(n²) cost",
    "",
    "Solution: Hybrid SSM + Local Attention architecture",
    "  → Block-wise Mamba scan balances temporal memory vs spatial coherence",
    "  → Frame local attention ensures short-term quality",
    "  → Clean-prefix training forces long-range dependency learning",
    "",
    "Results:",
    "  → Linear training complexity, constant inference cost",
    "  → Matches full-context transformers at fraction of the cost",
    "  → All three components proven essential via ablation",
    "",
    "Impact: First video world model with genuine long-term memory",
    "at practical computational cost",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0),
                 items, font_size=22, color=LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Save
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
output_path = "/Users/ana/Desktop/Timur/thesis_ve/reports/arxiv_2505_20171_presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
